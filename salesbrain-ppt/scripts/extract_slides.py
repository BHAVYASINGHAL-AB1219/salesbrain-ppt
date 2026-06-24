#!/usr/bin/env python3
"""
extract_slides.py — Sales Brain PPT Reviewer Utility

Reads a .pptx file using python-pptx and extracts a structured JSON
summary of each slide's visible text content, shapes, tables, charts,
AND shape geometry (positions, sizes) for visual/structural review.

Usage:
    python3 scripts/extract_slides.py /path/to/deck.pptx

Output:
    Prints JSON to stdout, consumed by the Node.js reviewerAgent.
"""

import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


# Standard 16:9 slide dimensions in inches
SLIDE_WIDTH_INCHES = 10.0
SLIDE_HEIGHT_INCHES = 5.625  # 7.5 * (9/12) for widescreen — actually 5.625 for pptxgenjs LAYOUT_16x9


def emu_to_inches(emu_value):
    """Convert EMU (English Metric Units) to inches."""
    if emu_value is None:
        return 0.0
    return round(emu_value / 914400, 4)


def get_shape_type_name(shape):
    """Get a human-readable shape type string."""
    try:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return "auto_shape"
        elif st == MSO_SHAPE_TYPE.TEXT_BOX:
            return "text_box"
        elif st == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        elif st == MSO_SHAPE_TYPE.CHART:
            return "chart"
        elif st == MSO_SHAPE_TYPE.TABLE:
            return "table"
        elif st == MSO_SHAPE_TYPE.GROUP:
            return "group"
        elif st == MSO_SHAPE_TYPE.FREEFORM:
            return "freeform"
        elif st == MSO_SHAPE_TYPE.PLACEHOLDER:
            return "placeholder"
        else:
            return str(st).split('.')[-1].lower() if st else "unknown"
    except Exception:
        return "unknown"


def get_min_font_size(shape):
    """Get the smallest font size (in pt) used in a shape's text frame."""
    if not shape.has_text_frame:
        return None
    
    min_size = None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                size_pt = run.font.size.pt
                if min_size is None or size_pt < min_size:
                    min_size = size_pt
    return min_size


def get_max_font_size(shape):
    """Get the largest font size (in pt) used in a shape's text frame."""
    if not shape.has_text_frame:
        return None
    
    max_size = None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                size_pt = run.font.size.pt
                if max_size is None or size_pt > max_size:
                    max_size = size_pt
    return max_size


def get_text_length(shape):
    """Get total character count of visible text in a shape."""
    if not shape.has_text_frame:
        return 0
    return sum(len(para.text.strip()) for para in shape.text_frame.paragraphs)


def get_autosize_mode(shape):
    """Get the text frame auto-size mode."""
    if not shape.has_text_frame:
        return None
    try:
        autosize = shape.text_frame.auto_size
        if autosize is None:
            return "NONE"
        return str(autosize).split('.')[-1]
    except Exception:
        return "UNKNOWN"


def extract_shape_geometry(shape, shape_index):
    """Extract position and size geometry for a shape."""
    return {
        "shape_index": shape_index,
        "shape_type": get_shape_type_name(shape),
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "right": emu_to_inches(shape.left + shape.width) if shape.left is not None and shape.width is not None else 0,
        "bottom": emu_to_inches(shape.top + shape.height) if shape.top is not None and shape.height is not None else 0,
        "has_text_frame": shape.has_text_frame,
        "text_length": get_text_length(shape),
        "font_size_min": get_min_font_size(shape),
        "font_size_max": get_max_font_size(shape),
        "autosize_mode": get_autosize_mode(shape),
        "name": shape.name or f"Shape_{shape_index}",
    }


def extract_text_from_shape(shape):
    """Extract all text from a shape's text frame, preserving paragraph structure."""
    if not shape.has_text_frame:
        return []

    paragraphs = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Detect style from font properties
        style = "body"
        if para.runs:
            run = para.runs[0]
            font = run.font
            if font.size and font.size >= Pt(28):
                style = "title"
            elif font.size and font.size >= Pt(20):
                style = "subtitle"
            elif font.size and font.size <= Pt(10):
                style = "eyebrow"
            if font.bold:
                style = "title" if font.size and font.size >= Pt(20) else "bold_" + style

        paragraphs.append({
            "text": text,
            "style": style
        })

    return paragraphs


def extract_table(shape):
    """Extract table data as a list of rows."""
    if not shape.has_table:
        return None

    table = shape.table
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(cell.text.strip())
        rows.append(cells)

    return {
        "type": "table",
        "rows": rows,
        "row_count": len(rows),
        "col_count": len(rows[0]) if rows else 0
    }


def extract_chart(shape):
    """Extract basic chart metadata."""
    if not shape.has_chart:
        return None

    chart = shape.chart
    chart_data = {
        "type": "chart",
        "chart_type": str(chart.chart_type),
        "has_legend": chart.has_legend,
    }

    # Try to extract series data
    try:
        series_info = []
        for series in chart.series:
            series_info.append({
                "name": str(series.name) if series.name else "Unnamed",
                "value_count": len(series.values) if hasattr(series, 'values') else 0
            })
        chart_data["series"] = series_info
    except Exception:
        chart_data["series"] = []

    return chart_data


def extract_slide(slide, slide_index, slide_width, slide_height):
    """Extract all content from a single slide, including shape geometry."""
    text_blocks = []
    tables = []
    charts = []
    shapes_geometry = []
    image_count = 0
    shape_count = len(slide.shapes)

    for idx, shape in enumerate(slide.shapes):
        # Extract geometry for EVERY shape
        geo = extract_shape_geometry(shape, idx)
        shapes_geometry.append(geo)

        # Extract text
        text_paragraphs = extract_text_from_shape(shape)
        if text_paragraphs:
            text_blocks.extend(text_paragraphs)

        # Extract tables
        if shape.has_table:
            table_data = extract_table(shape)
            if table_data:
                tables.append(table_data)

        # Extract charts
        if shape.has_chart:
            chart_data = extract_chart(shape)
            if chart_data:
                charts.append(chart_data)

        # Count images
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_count += 1

    # Get slide notes if present
    notes_text = ""
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame:
            notes_text = notes_slide.notes_text_frame.text.strip()

    return {
        "slide_number": slide_index + 1,
        "text_blocks": text_blocks,
        "all_text_combined": " | ".join([tb["text"] for tb in text_blocks]),
        "tables": tables,
        "charts": charts,
        "shapes_geometry": shapes_geometry,
        "image_count": image_count,
        "shape_count": shape_count,
        "speaker_notes": notes_text,
        "has_chart": len(charts) > 0,
        "has_table": len(tables) > 0,
        "has_images": image_count > 0,
        "text_block_count": len(text_blocks)
    }


def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: python3 extract_slides.py <path_to_pptx>"}))
        sys.exit(1)

    pptx_path = sys.argv[1]

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(json.dumps({"error": f"Failed to open {pptx_path}: {str(e)}"}))
        sys.exit(1)

    # Get slide dimensions
    slide_width = emu_to_inches(prs.slide_width)
    slide_height = emu_to_inches(prs.slide_height)

    slides_data = []
    for i, slide in enumerate(prs.slides):
        slide_data = extract_slide(slide, i, slide_width, slide_height)
        slides_data.append(slide_data)

    output = {
        "file_path": pptx_path,
        "slide_count": len(slides_data),
        "slide_width_inches": slide_width,
        "slide_height_inches": slide_height,
        "slides": slides_data
    }

    print(json.dumps(output, ensure_ascii=False))


if __name__ == "__main__":
    main()
